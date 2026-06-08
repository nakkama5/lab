"""Document ingestion for .pptx, .pdf, .docx, .txt, .md files."""
from __future__ import annotations

import os
from pathlib import Path
from typing import Any


def _ingest_pptx(path: str) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(path)
    chunks = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
        # Slide notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                text_parts.append(f"[Notes] {notes_text}")
        if text_parts:
            chunks.append({
                "source_file": os.path.basename(path),
                "section": f"Slide {slide_num}",
                "text": "\n".join(text_parts),
            })
    return chunks


def _ingest_pdf(path: str) -> list[dict]:
    import pdfplumber

    chunks = []
    with pdfplumber.open(path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            text = page.extract_text()
            if text and text.strip():
                chunks.append({
                    "source_file": os.path.basename(path),
                    "section": f"Page {page_num}",
                    "text": text.strip(),
                })
    return chunks


def _ingest_docx(path: str) -> list[dict]:
    from docx import Document

    doc = Document(path)
    chunks = []
    current_section = "Document"
    current_parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # Use headings as section delimiters
        if para.style.name.startswith("Heading"):
            if current_parts:
                chunks.append({
                    "source_file": os.path.basename(path),
                    "section": current_section,
                    "text": "\n".join(current_parts),
                })
                current_parts = []
            current_section = text
        else:
            current_parts.append(text)

    if current_parts:
        chunks.append({
            "source_file": os.path.basename(path),
            "section": current_section,
            "text": "\n".join(current_parts),
        })
    return chunks


def _ingest_text(path: str) -> list[dict]:
    content = Path(path).read_text(encoding="utf-8", errors="replace").strip()
    if not content:
        return []
    return [{
        "source_file": os.path.basename(path),
        "section": "Full document",
        "text": content,
    }]


def ingest_files(paths: list[str]) -> list[dict]:
    """Ingest a list of file paths into a list of text chunks.

    Each chunk: {source_file, section, text}
    """
    chunks: list[dict] = []
    for path in paths:
        ext = Path(path).suffix.lower()
        try:
            if ext == ".pptx":
                chunks.extend(_ingest_pptx(path))
            elif ext == ".pdf":
                chunks.extend(_ingest_pdf(path))
            elif ext == ".docx":
                chunks.extend(_ingest_docx(path))
            elif ext in (".txt", ".md"):
                chunks.extend(_ingest_text(path))
            else:
                # Try plain text as fallback
                chunks.extend(_ingest_text(path))
        except Exception as e:
            chunks.append({
                "source_file": os.path.basename(path),
                "section": "ERROR",
                "text": f"Failed to ingest: {e}",
            })
    return chunks


def load_files(uploaded_files) -> list[dict]:
    """Ingest Streamlit UploadedFile objects. Returns list of text chunks."""
    import tempfile
    chunks: list[dict] = []
    for uf in uploaded_files:
        suffix = Path(uf.name).suffix.lower()
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(uf.read())
                tmp_path = tmp.name
            chunks.extend(ingest_files([tmp_path]))
            # Fix source_file name
            for c in chunks:
                if c.get("source_file", "").startswith("tmp"):
                    c["source_file"] = uf.name
            Path(tmp_path).unlink(missing_ok=True)
        except Exception as e:
            chunks.append({"source_file": uf.name, "section": "ERROR", "text": str(e)})
    return chunks
    """Flatten corpus chunks to a single text string for LLM consumption."""
    parts = []
    for chunk in corpus:
        parts.append(
            f"[{chunk['source_file']} / {chunk['section']}]\n{chunk['text']}"
        )
    return "\n\n---\n\n".join(parts)
