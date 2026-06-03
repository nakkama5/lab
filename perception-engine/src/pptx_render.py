"""Full python-pptx renderer for all 14 section types."""
from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu


# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

FOOTER_H = Inches(0.3)
FOOTER_Y = SLIDE_H - FOOTER_H


# ── Color helpers ─────────────────────────────────────────────────────────────

def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert a 6-digit hex string (no #) to RGBColor."""
    h = hex_str.lstrip("#")
    if len(h) != 6:
        h = (h + "000000")[:6]
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)


def _palette(brand_tokens: dict) -> dict[str, RGBColor]:
    """Extract palette as a dict of name -> RGBColor."""
    palette = brand_tokens.get("palette", {})
    result: dict[str, RGBColor] = {}
    for name, info in palette.items():
        result[name] = hex_to_rgb(info.get("hex", "000000"))
    # Ensure defaults exist
    defaults = {
        "ink": "17120E", "jade": "2C5D4F", "saffron": "E1A23C",
        "clay": "B5683E", "cream": "F4EEE2",
    }
    for name, hex_val in defaults.items():
        if name not in result:
            result[name] = hex_to_rgb(hex_val)
    return result


def _font_names(brand_tokens: dict) -> tuple[str, str]:
    """Return (display_font, body_font)."""
    type_info = brand_tokens.get("type", {})
    return type_info.get("display", "Georgia"), type_info.get("body", "Calibri")


# ── Shape helpers ─────────────────────────────────────────────────────────────

def _add_rect(slide, x, y, w, h, fill_color: RGBColor | None = None) -> Any:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h,
    )
    shape.line.fill.background()
    shape.line.width = Pt(0)
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def _add_textbox(
    slide, text: str, x, y, w, h,
    font_name: str = "Calibri",
    font_size: float = 12,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
) -> Any:
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def _add_footer(slide, page_num: int, colors: dict[str, RGBColor], body_font: str) -> None:
    """Add confidential footer and page number to every slide."""
    cream = colors["cream"]
    # Left: CONFIDENTIAL text
    _add_textbox(
        slide, "CONFIDENTIAL — INTERNAL USE ONLY",
        Inches(0.2), FOOTER_Y, Inches(6), FOOTER_H,
        font_name=body_font, font_size=7, color=cream,
        align=PP_ALIGN.LEFT,
    )
    # Right: page number
    _add_textbox(
        slide, str(page_num),
        Inches(12.5), FOOTER_Y, Inches(0.6), FOOTER_H,
        font_name=body_font, font_size=7, color=cream,
        align=PP_ALIGN.RIGHT,
    )


# ── Section renderers ─────────────────────────────────────────────────────────

def render_cover(
    prs: Presentation,
    deck_spec: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
) -> None:
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark ink background
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["ink"])

    product = deck_spec.get("product", "")
    taglines = deck_spec.get("taglines", {})
    punchy = taglines.get("punchy", "")
    micro = deck_spec.get("micro", "")

    # Product name — large
    _add_textbox(
        slide, product,
        Inches(1.0), Inches(1.5), Inches(11.33), Inches(2.0),
        font_name=display_font, font_size=60, bold=True,
        color=colors["cream"], align=PP_ALIGN.LEFT,
    )
    # Punchy tagline
    if punchy:
        _add_textbox(
            slide, punchy,
            Inches(1.0), Inches(3.7), Inches(11.33), Inches(1.0),
            font_name=display_font, font_size=22,
            color=colors["saffron"], align=PP_ALIGN.LEFT,
        )
    # Micro copy
    if micro:
        _add_textbox(
            slide, micro,
            Inches(1.0), Inches(5.0), Inches(10.0), Inches(0.8),
            font_name=body_font, font_size=12,
            color=colors["cream"], align=PP_ALIGN.LEFT,
        )

    _add_footer(slide, page_num, colors, body_font)


def render_metrics(
    prs: Presentation,
    deck_spec: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
    count: int = 4,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Cream background
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["cream"])

    # Title
    _add_textbox(
        slide, "The Numbers",
        Inches(0.5), Inches(0.3), Inches(12.0), Inches(0.7),
        font_name=display_font, font_size=28, bold=True,
        color=colors["ink"], align=PP_ALIGN.LEFT,
    )

    metrics = deck_spec.get("metrics", [])[:count]
    n = max(len(metrics), 1)
    box_w = Inches(12.0 / n)
    box_h = Inches(4.0)
    top_y = Inches(1.5)

    for i, m in enumerate(metrics):
        x = Inches(0.5) + i * box_w
        # Jade box background
        _add_rect(slide, x + Inches(0.1), top_y, box_w - Inches(0.2), box_h, colors["jade"])
        # Large number in saffron
        _add_textbox(
            slide, m.get("num", ""),
            x + Inches(0.15), top_y + Inches(0.5), box_w - Inches(0.3), Inches(1.8),
            font_name=display_font, font_size=48, bold=True,
            color=colors["saffron"], align=PP_ALIGN.CENTER,
        )
        # Label below
        _add_textbox(
            slide, m.get("label", ""),
            x + Inches(0.15), top_y + Inches(2.4), box_w - Inches(0.3), Inches(1.2),
            font_name=body_font, font_size=14,
            color=colors["cream"], align=PP_ALIGN.CENTER,
        )

    _add_footer(slide, page_num, colors, body_font)


def render_legacy_evolution(
    prs: Presentation,
    deck_spec: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
    count: int = 5,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["cream"])

    # Left/right headers
    col_w = Inches(6.0)
    _add_rect(slide, Inches(0.3), Inches(0.2), col_w, Inches(0.8), colors["clay"])
    _add_textbox(
        slide, "Legacy", Inches(0.3), Inches(0.2), col_w, Inches(0.8),
        font_name=display_font, font_size=22, bold=True,
        color=colors["cream"], align=PP_ALIGN.CENTER,
    )
    _add_rect(slide, Inches(6.8), Inches(0.2), col_w, Inches(0.8), colors["jade"])
    _add_textbox(
        slide, "Evolution", Inches(6.8), Inches(0.2), col_w, Inches(0.8),
        font_name=display_font, font_size=22, bold=True,
        color=colors["cream"], align=PP_ALIGN.CENTER,
    )

    legacy = deck_spec.get("legacy", [])
    evolution = deck_spec.get("evolution", [])
    row_h = Inches(5.8 / count)

    for i in range(count):
        y = Inches(1.2) + i * row_h
        legacy_text = legacy[i] if i < len(legacy) else ""
        evol_text = evolution[i] if i < len(evolution) else ""
        # Alternating row bg
        if i % 2 == 0:
            _add_rect(slide, Inches(0.3), y, col_w, row_h, hex_to_rgb("EDE7DA"))
            _add_rect(slide, Inches(6.8), y, col_w, row_h, hex_to_rgb("D6E8E3"))
        _add_textbox(
            slide, legacy_text, Inches(0.5), y + Inches(0.05), col_w - Inches(0.4), row_h - Inches(0.1),
            font_name=body_font, font_size=13, color=colors["ink"],
        )
        _add_textbox(
            slide, evol_text, Inches(7.0), y + Inches(0.05), col_w - Inches(0.4), row_h - Inches(0.1),
            font_name=body_font, font_size=13, color=colors["ink"],
        )

    _add_footer(slide, page_num, colors, body_font)


def render_jargon_to_value(
    prs: Presentation,
    deck_spec: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
    count: int = 4,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["cream"])

    headers = ["Feature", "Capability", "Benefit", "KPI"]
    col_w = Inches(12.0 / 4)

    # Header row
    for j, header in enumerate(headers):
        x = Inches(0.5) + j * col_w
        _add_rect(slide, x, Inches(0.2), col_w - Inches(0.05), Inches(0.7), colors["jade"])
        _add_textbox(
            slide, header, x, Inches(0.2), col_w - Inches(0.05), Inches(0.7),
            font_name=display_font, font_size=14, bold=True,
            color=colors["cream"], align=PP_ALIGN.CENTER,
        )

    rows = deck_spec.get("jargon_rows", [])[:count]
    row_h = Inches(5.8 / max(len(rows), 1))

    for i, row in enumerate(rows):
        values = [
            row.get("feature", ""), row.get("capability", ""),
            row.get("benefit", ""), row.get("kpi", ""),
        ]
        y = Inches(1.1) + i * row_h
        bg = hex_to_rgb("EDE7DA") if i % 2 == 0 else colors["cream"]
        for j, val in enumerate(values):
            x = Inches(0.5) + j * col_w
            _add_rect(slide, x, y, col_w - Inches(0.05), row_h - Inches(0.05), bg)
            _add_textbox(
                slide, val, x + Inches(0.05), y + Inches(0.05),
                col_w - Inches(0.15), row_h - Inches(0.1),
                font_name=body_font, font_size=12, color=colors["ink"],
            )

    _add_footer(slide, page_num, colors, body_font)


def render_swot(
    prs: Presentation,
    deck_spec: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["cream"])

    swot = deck_spec.get("swot", {})
    quadrants = [
        ("Strengths", swot.get("strengths", []), colors["jade"], 0, 0),
        ("Weaknesses", swot.get("weaknesses", []), colors["clay"], 1, 0),
        ("Opportunities", swot.get("opportunities", []), hex_to_rgb("4A7C6F"), 0, 1),
        ("Threats", swot.get("threats", []), hex_to_rgb("8B4513"), 1, 1),
    ]

    quad_w = Inches(6.4)
    quad_h = Inches(3.4)

    for label, items, header_color, col, row in quadrants:
        x = Inches(0.3) + col * quad_w
        y = Inches(0.3) + row * quad_h

        # Ink header strip
        _add_rect(slide, x, y, quad_w - Inches(0.1), Inches(0.55), header_color)
        _add_textbox(
            slide, label, x, y, quad_w - Inches(0.1), Inches(0.55),
            font_name=display_font, font_size=16, bold=True,
            color=colors["cream"], align=PP_ALIGN.CENTER,
        )
        # Content area
        _add_rect(slide, x, y + Inches(0.55), quad_w - Inches(0.1), quad_h - Inches(0.65), hex_to_rgb("FAF6EF"))
        bullets = "\n".join(f"• {item}" for item in items)
        _add_textbox(
            slide, bullets,
            x + Inches(0.1), y + Inches(0.65),
            quad_w - Inches(0.3), quad_h - Inches(0.8),
            font_name=body_font, font_size=11, color=colors["ink"],
        )

    _add_footer(slide, page_num, colors, body_font)


def render_metaphor(
    prs: Presentation,
    deck_spec: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["jade"])

    metaphor = deck_spec.get("metaphor", {})
    statement = metaphor.get("statement", "")
    rationale = metaphor.get("rationale", "")

    _add_textbox(
        slide, statement,
        Inches(1.0), Inches(1.5), Inches(11.33), Inches(2.5),
        font_name=display_font, font_size=36, bold=True,
        color=colors["cream"], align=PP_ALIGN.CENTER,
    )
    if rationale:
        _add_textbox(
            slide, rationale,
            Inches(1.5), Inches(4.5), Inches(10.33), Inches(2.0),
            font_name=body_font, font_size=14,
            color=colors["cream"], align=PP_ALIGN.CENTER,
        )

    _add_footer(slide, page_num, colors, body_font)


def render_palette(
    prs: Presentation,
    brand_tokens: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["cream"])

    _add_textbox(
        slide, "Colour Palette",
        Inches(0.5), Inches(0.2), Inches(10), Inches(0.7),
        font_name=display_font, font_size=28, bold=True,
        color=colors["ink"],
    )

    palette = brand_tokens.get("palette", {})
    items = list(palette.items())
    n = len(items)
    swatch_w = Inches(12.0 / max(n, 1))
    swatch_h = Inches(3.5)
    top_y = Inches(1.2)

    for i, (name, info) in enumerate(items):
        x = Inches(0.5) + i * swatch_w
        color = hex_to_rgb(info.get("hex", "888888"))
        _add_rect(slide, x, top_y, swatch_w - Inches(0.1), swatch_h, color)
        # Color info below swatch
        _add_textbox(
            slide, f"{name.capitalize()}\n#{info.get('hex', '')}\n{info.get('role', '')}",
            x, top_y + swatch_h + Inches(0.05), swatch_w - Inches(0.1), Inches(1.5),
            font_name=body_font, font_size=10, color=colors["ink"], align=PP_ALIGN.CENTER,
        )

    _add_footer(slide, page_num, colors, body_font)


def render_tone_pillars(
    prs: Presentation,
    deck_spec: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
    count: int = 3,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["cream"])

    _add_textbox(
        slide, "Tone of Voice",
        Inches(0.5), Inches(0.2), Inches(10), Inches(0.7),
        font_name=display_font, font_size=28, bold=True,
        color=colors["ink"],
    )

    pillars = deck_spec.get("pillars", [])[:count]
    col_w = Inches(12.0 / max(len(pillars), 1))

    for i, pillar in enumerate(pillars):
        x = Inches(0.5) + i * col_w
        _add_rect(slide, x, Inches(1.1), col_w - Inches(0.2), Inches(0.6), colors["jade"])
        _add_textbox(
            slide, pillar.get("name", ""),
            x, Inches(1.1), col_w - Inches(0.2), Inches(0.6),
            font_name=display_font, font_size=16, bold=True,
            color=colors["cream"], align=PP_ALIGN.CENTER,
        )
        _add_textbox(
            slide, f"✓ Do say:\n{pillar.get('do_say', '')}",
            x + Inches(0.1), Inches(1.9), col_w - Inches(0.3), Inches(2.0),
            font_name=body_font, font_size=12, color=colors["ink"],
        )
        _add_textbox(
            slide, f"✗ Don't say:\n{pillar.get('dont_say', '')}",
            x + Inches(0.1), Inches(4.2), col_w - Inches(0.3), Inches(2.0),
            font_name=body_font, font_size=12, color=colors["clay"],
        )

    _add_footer(slide, page_num, colors, body_font)


def render_vocab(
    prs: Presentation,
    deck_spec: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
    count: int = 6,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["cream"])

    _add_textbox(
        slide, "Vocabulary Swap",
        Inches(0.5), Inches(0.2), Inches(10), Inches(0.7),
        font_name=display_font, font_size=28, bold=True,
        color=colors["ink"],
    )

    # Headers
    col_w = Inches(5.5)
    _add_rect(slide, Inches(0.5), Inches(1.0), col_w, Inches(0.5), colors["clay"])
    _add_textbox(
        slide, "Say goodbye to…", Inches(0.5), Inches(1.0), col_w, Inches(0.5),
        font_name=display_font, font_size=14, bold=True,
        color=colors["cream"], align=PP_ALIGN.CENTER,
    )
    _add_rect(slide, Inches(7.0), Inches(1.0), col_w, Inches(0.5), colors["jade"])
    _add_textbox(
        slide, "Say hello to…", Inches(7.0), Inches(1.0), col_w, Inches(0.5),
        font_name=display_font, font_size=14, bold=True,
        color=colors["cream"], align=PP_ALIGN.CENTER,
    )

    vocab = deck_spec.get("vocab", [])[:count]
    row_h = Inches(5.5 / max(len(vocab), 1))

    for i, item in enumerate(vocab):
        y = Inches(1.6) + i * row_h
        from_text = item.get("from", item.get("from_", ""))
        to_text = item.get("to", "")
        bg = hex_to_rgb("EDE7DA") if i % 2 == 0 else colors["cream"]
        _add_rect(slide, Inches(0.5), y, col_w, row_h - Inches(0.05), bg)
        _add_textbox(
            slide, from_text, Inches(0.6), y + Inches(0.05), col_w - Inches(0.2), row_h - Inches(0.1),
            font_name=body_font, font_size=13, color=colors["clay"],
        )
        _add_rect(slide, Inches(7.0), y, col_w, row_h - Inches(0.05), hex_to_rgb("D6E8E3") if i % 2 == 0 else colors["cream"])
        _add_textbox(
            slide, to_text, Inches(7.1), y + Inches(0.05), col_w - Inches(0.2), row_h - Inches(0.1),
            font_name=body_font, font_size=13, color=colors["jade"],
        )

    _add_footer(slide, page_num, colors, body_font)


def render_taglines(
    prs: Presentation,
    deck_spec: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["ink"])

    _add_textbox(
        slide, "The Baselines",
        Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
        font_name=display_font, font_size=28, bold=True,
        color=colors["cream"],
    )

    taglines = deck_spec.get("taglines", {})
    variants = [
        ("Outcome-led", taglines.get("outcome", "")),
        ("Visionary", taglines.get("visionary", "")),
        ("Punchy", taglines.get("punchy", "")),
    ]

    for i, (label, text) in enumerate(variants):
        y = Inches(1.4) + i * Inches(1.8)
        _add_textbox(
            slide, label.upper(),
            Inches(0.8), y, Inches(3.0), Inches(0.5),
            font_name=body_font, font_size=10, bold=True,
            color=colors["saffron"],
        )
        _add_textbox(
            slide, text,
            Inches(0.8), y + Inches(0.5), Inches(11.5), Inches(1.0),
            font_name=display_font, font_size=24, bold=True,
            color=colors["cream"],
        )

    _add_footer(slide, page_num, colors, body_font)


def render_manifesto(
    prs: Presentation,
    deck_spec: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["jade"])

    manifesto = deck_spec.get("manifesto", "")
    _add_textbox(
        slide, manifesto,
        Inches(1.0), Inches(0.8), Inches(11.33), Inches(5.8),
        font_name=display_font, font_size=18, bold=False,
        color=colors["cream"], align=PP_ALIGN.CENTER,
    )

    _add_footer(slide, page_num, colors, body_font)


def render_roadmap(
    prs: Presentation,
    deck_spec: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
    count: int = 3,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["cream"])

    _add_textbox(
        slide, "Roadmap",
        Inches(0.5), Inches(0.2), Inches(10), Inches(0.7),
        font_name=display_font, font_size=28, bold=True,
        color=colors["ink"],
    )

    phases = deck_spec.get("roadmap", [])[:count]
    box_w = Inches(12.0 / max(len(phases), 1))
    box_h = Inches(5.5)
    top_y = Inches(1.1)

    phase_colors = [colors["jade"], colors["clay"], colors["saffron"]]

    for i, phase in enumerate(phases):
        x = Inches(0.5) + i * box_w
        pc = phase_colors[i % len(phase_colors)]
        _add_rect(slide, x, top_y, box_w - Inches(0.2), Inches(0.7), pc)
        _add_textbox(
            slide, f"Phase {phase.get('phase', i+1)}",
            x, top_y, box_w - Inches(0.2), Inches(0.7),
            font_name=display_font, font_size=14, bold=True,
            color=colors["cream"], align=PP_ALIGN.CENTER,
        )
        _add_rect(slide, x, top_y + Inches(0.7), box_w - Inches(0.2), box_h - Inches(0.7), hex_to_rgb("FAF6EF"))
        phase_name = phase.get("name", "")
        when = phase.get("when", "")
        points = phase.get("points", [])
        content = f"{phase_name}\n{when}\n\n" + "\n".join(f"• {p}" for p in points if p)
        _add_textbox(
            slide, content,
            x + Inches(0.1), top_y + Inches(0.8), box_w - Inches(0.4), box_h - Inches(0.9),
            font_name=body_font, font_size=12, color=colors["ink"],
        )

    _add_footer(slide, page_num, colors, body_font)


def render_grapevine(
    prs: Presentation,
    deck_spec: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
    count: int = 3,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["cream"])

    _add_textbox(
        slide, "Bypassing the Critics",
        Inches(0.5), Inches(0.2), Inches(10), Inches(0.7),
        font_name=display_font, font_size=28, bold=True,
        color=colors["ink"],
    )

    items = deck_spec.get("grapevine", [])[:count]
    card_w = Inches(12.0 / max(len(items), 1))
    card_h = Inches(5.0)
    top_y = Inches(1.2)

    for i, item in enumerate(items):
        x = Inches(0.5) + i * card_w
        _add_rect(slide, x, top_y, card_w - Inches(0.2), card_h, colors["jade"])
        _add_textbox(
            slide, item.get("title", ""),
            x + Inches(0.1), top_y + Inches(0.2), card_w - Inches(0.4), Inches(0.8),
            font_name=display_font, font_size=16, bold=True,
            color=colors["saffron"],
        )
        _add_textbox(
            slide, item.get("desc", ""),
            x + Inches(0.1), top_y + Inches(1.1), card_w - Inches(0.4), card_h - Inches(1.4),
            font_name=body_font, font_size=12, color=colors["cream"],
        )

    _add_footer(slide, page_num, colors, body_font)


def render_closing(
    prs: Presentation,
    deck_spec: dict,
    colors: dict,
    display_font: str,
    body_font: str,
    page_num: int,
) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, colors["ink"])

    product = deck_spec.get("product", "")
    taglines = deck_spec.get("taglines", {})
    punchy = taglines.get("punchy", "")

    _add_textbox(
        slide, product,
        Inches(1.0), Inches(2.5), Inches(11.33), Inches(1.5),
        font_name=display_font, font_size=52, bold=True,
        color=colors["cream"], align=PP_ALIGN.CENTER,
    )
    if punchy:
        _add_textbox(
            slide, punchy,
            Inches(1.0), Inches(4.3), Inches(11.33), Inches(1.0),
            font_name=display_font, font_size=20,
            color=colors["saffron"], align=PP_ALIGN.CENTER,
        )

    _add_footer(slide, page_num, colors, body_font)


# ── Main render entry point ───────────────────────────────────────────────────

def render_deck(
    deck_spec: dict,
    brand_tokens: dict,
    deck_layout: dict,
    output_path: str,
) -> None:
    """Render all enabled sections to a PPTX file."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    colors = _palette(brand_tokens)
    display_font, body_font = _font_names(brand_tokens)

    sections = deck_layout.get("sections", [])
    enabled = [s for s in sections if s.get("enabled", True)]

    page_num = 1

    for section in enabled:
        stype = section.get("type", "")
        count = section.get("count")

        kwargs: dict = dict(
            prs=prs,
            colors=colors,
            display_font=display_font,
            body_font=body_font,
            page_num=page_num,
        )

        if stype == "cover":
            render_cover(**kwargs, deck_spec=deck_spec)
        elif stype == "metrics":
            render_metrics(**kwargs, deck_spec=deck_spec, count=count or 4)
        elif stype == "legacy_evolution":
            render_legacy_evolution(**kwargs, deck_spec=deck_spec, count=count or 5)
        elif stype == "jargon_to_value":
            render_jargon_to_value(**kwargs, deck_spec=deck_spec, count=count or 4)
        elif stype == "swot":
            render_swot(**kwargs, deck_spec=deck_spec)
        elif stype == "metaphor":
            render_metaphor(**kwargs, deck_spec=deck_spec)
        elif stype == "palette":
            render_palette(**kwargs, brand_tokens=brand_tokens)
        elif stype == "tone_pillars":
            render_tone_pillars(**kwargs, deck_spec=deck_spec, count=count or 3)
        elif stype == "vocab":
            render_vocab(**kwargs, deck_spec=deck_spec, count=count or 6)
        elif stype == "taglines":
            render_taglines(**kwargs, deck_spec=deck_spec)
        elif stype == "manifesto":
            render_manifesto(**kwargs, deck_spec=deck_spec)
        elif stype == "roadmap":
            render_roadmap(**kwargs, deck_spec=deck_spec, count=count or 3)
        elif stype == "grapevine":
            render_grapevine(**kwargs, deck_spec=deck_spec, count=count or 3)
        elif stype == "closing":
            render_closing(**kwargs, deck_spec=deck_spec)
        else:
            # Unknown section type — skip
            continue

        page_num += 1

    prs.save(output_path)
